import json
import sys

from pptx import Presentation


def main():
    data = json.loads(sys.argv[1])
    output_path = sys.argv[2]

    presentation = Presentation()

    for slide_data in data["slides"]:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_data["title"]

        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        for index, item in enumerate(slide_data["content"]):
            if index == 0:
                text_frame.text = item
            else:
                text_frame.add_paragraph().text = item

    presentation.save(output_path)


main()
